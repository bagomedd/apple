from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide Layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Early Personal Computers"
subtitle.text = "Created by [Your Name]"

# Slide 2: Altair by MITS (1974)
slide_layout = prs.slide_layouts[5]  # Title Only Layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Altair by MITS (1974)"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
text_frame = content.text_frame
text_frame.text = "The MITS Altair 8800 is considered the first successful personal computer. It featured an Intel 8080 CPU and was sold as a kit."

# Slide 3: Tandy by Tandy Corporation (1977)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Tandy by Tandy Corporation (1977)"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
text_frame = content.text_frame
text_frame.text = "The Tandy TRS-80 was an early mass-market personal computer, known for its affordability and availability in RadioShack stores."

# Slide 4: Apple Macintosh (1978)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Apple Macintosh (1978)"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
text_frame = content.text_frame
text_frame.text = "The Apple Macintosh was one of the first computers to feature a graphical user interface (GUI), revolutionizing personal computing."

# Slide 5: IBM PC (1980)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "IBM PC (1980)"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
text_frame = content.text_frame
text_frame.text = "The IBM PC became the industry standard for personal computers, using an open architecture that led to widespread adoption."

# Slide 6: Resources Slide
slide_layout = prs.slide_layouts[5]  # Title Only Layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Resources"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
text_frame = content.text_frame
text_frame.text = "Sources:\n- Wikipedia\n- Computer History Museum\n- Various online archives"

# Save the presentation
pptx_path = "/mnt/data/Early_Personal_Computers.pptx"
prs.save(pptx_path)

pptx_path
